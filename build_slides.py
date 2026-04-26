# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx>=1.0"]
# ///
"""Generate slides.pptx mirroring slides.html content.

Run with:
    uv run build_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x0C, 0x23, 0x4B)
NAVY_DEEP = RGBColor(0x06, 0x12, 0x2A)
RED = RGBColor(0xC8, 0x10, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xC9, 0xCF, 0xDC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.shadow.inherit = False
    return shape


def _add_text(slide, left, top, width, height, text, *, size, color=WHITE,
              bold=False, italic=False, tracking=0, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _add_eyebrow(slide, text):
    _add_text(slide, Inches(0.9), Inches(0.7), Inches(10), Inches(0.4),
              text.upper(), size=14, color=RED, bold=True)


def _add_slide_num(slide, n):
    _add_text(slide, Inches(11.4), Inches(0.35), Inches(1.7), Inches(0.4),
              f"{n:02d} / 05", size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def _add_h2(slide, text, top=Inches(1.15)):
    _add_text(slide, Inches(0.9), top, Inches(11.5), Inches(0.9),
              text, size=36, color=WHITE, bold=True)


def _add_bullets(slide, items, top=Inches(2.4), left=Inches(0.9),
                 width=Inches(11.5), height=Inches(4.7), size=18):
    """items: list of (plain_text, [(text, bold), ...])

    If the list-item form is the second element, runs get individual bold flags.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0

    # hanging-indent width tuned so wrapped lines align with the text after the dash
    indent_emu = int(0.42 * 914400)

    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        p.space_before = Pt(0)
        # set hanging indent: marL pushes every line right, indent pulls the first line back
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        pPr.set("marL", str(indent_emu))
        pPr.set("indent", str(-indent_emu))
        # red dash
        dash = p.add_run()
        dash.text = "—  "
        dash.font.name = "Helvetica Neue"
        dash.font.size = Pt(size)
        dash.font.bold = True
        dash.font.color.rgb = RED

        if isinstance(item, str):
            runs = [(item, False)]
        else:
            runs = item
        for text, bold in runs:
            r = p.add_run()
            r.text = text
            r.font.name = "Helvetica Neue"
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = WHITE


# -------- Build deck --------

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# Slide 1 — title
s1 = prs.slides.add_slide(blank)
_set_bg(s1, NAVY_DEEP)
# accent block
_add_rect(s1, Inches(0.9), Inches(2.0), Inches(0.18), Inches(3.4), RED)
_add_slide_num(s1, 1)
_add_text(s1, Inches(1.3), Inches(2.0), Inches(11), Inches(0.5),
          "HACKARIZONA — TEAM 10", size=14, color=RED, bold=True)
_add_text(s1, Inches(1.3), Inches(2.55), Inches(11), Inches(1.2),
          "Nuclear", size=66, color=WHITE, bold=True)
_add_text(s1, Inches(1.3), Inches(3.55), Inches(11), Inches(1.2),
          "Derating", size=66, color=RED, bold=True)
_add_text(s1, Inches(1.3), Inches(4.55), Inches(11), Inches(1.2),
          "Forecaster", size=66, color=WHITE, bold=True)
_add_text(s1, Inches(1.3), Inches(5.85), Inches(11), Inches(1.0),
          "Forecasting weather-driven nuclear power loss 1 to 14 days ahead so "
          "operators can capture energy that would otherwise be left on the table.",
          size=18, color=MUTED)

# Slide 2 — what we built
s2 = prs.slides.add_slide(blank)
_set_bg(s2, NAVY_DEEP)
_add_slide_num(s2, 2)
_add_eyebrow(s2, "The Project")
_add_h2(s2, "What We Built")
_add_bullets(s2, [
    [("Uses historical data to predict how much power a nuclear plant will "
      "actually be able to produce over the next ", False),
     ("14", True), (" days", False)],
    "Combines weather, river temperature, and streamflow with the plant's own historical performance",
    "Delivered through a simple web dashboard with a live US map and per-plant detail pages",
    "Lets operators see trouble coming before it shows up on the grid",
    "Built on public data only — no proprietary feeds, fully reproducible",
])

# Slide 3 — derating problem
s3 = prs.slides.add_slide(blank)
_set_bg(s3, NAVY_DEEP)
_add_slide_num(s3, 3)
_add_eyebrow(s3, "The Problem")
_add_h2(s3, "What Is Derating, And Why It Hurts")
_add_bullets(s3, [
    "Derating is when a plant is forced to run below full power — capacity is there, but it cannot be used",
    "Driven by hot weather and warm river water — cooling intake limits force operators to throttle down",
    "Hits hardest during heatwaves, exactly when the grid needs every megawatt the most",
    [("A single summer of derating can cost ", False), ("$50M+", True),
     (" at one plant, ", False), ("$200M+", True), (" industry-wide each year", False)],
    "Today operators react after the dip starts — there is no public forecast of when it is coming",
    "Result: wasted clean energy, expensive last-minute fossil backup, and avoidable grid strain",
], size=16)

# Slide 4 — how
s4 = prs.slides.add_slide(blank)
_set_bg(s4, NAVY_DEEP)
_add_slide_num(s4, 4)
_add_eyebrow(s4, "Our Approach")
_add_h2(s4, "How The Forecaster Solves It")
_add_bullets(s4, [
    "Outputs a three-tier alert for each day in the forecast window — operational, watch, or alert — "
    "alongside the predicted output percentage and an uncertainty band so operators see both the expected "
    "value and the range of likely outcomes",
    [("Trained on ", False), ("20+", True), (" years of plant power records, weather, and river data", False)],
    [("One model per forecast horizon (", False), ("1", True), (" through ", False), ("14", True),
     (" days) — every day on the curve gets its own prediction", False)],
    "Uses live weather forecasts to project conditions forward, then translates those into expected plant output",
    "Includes a replay mode proving the model would have called past heat-driven dips in advance",
    "Architecture is plant-agnostic — adding a new reactor is a single config entry",
], size=15)

# Slide 5 — impact (two columns)
s5 = prs.slides.add_slide(blank)
_set_bg(s5, NAVY_DEEP)
_add_slide_num(s5, 5)
_add_eyebrow(s5, "The Impact")
_add_h2(s5, "Who Wins When We Get This Right")

# Energy Optimization column
_add_text(s5, Inches(0.9), Inches(2.3), Inches(5.7), Inches(0.4),
          "ENERGY OPTIMIZATION", size=12, color=RED, bold=True)
_add_bullets(s5, [
    [("Recover lost megawatts by scheduling maintenance and refueling around predicted dips — "
      "capturing even ", False), ("10%", True),
     (" of the U.S. fleet's annual weather-driven derating losses is an estimated ", False),
     ("~$20M", True), (" per year", False)],
    "Smarter cooling tower usage — operators know when to spin them up and when they will not be needed, "
    "saving water and auxiliary power",
    "Plan backup power in advance — grid operators can line up cheaper replacement generation days ahead "
    "instead of paying premium rates for emergency peaker plants once a heatwave hits",
    "Lower carbon grid response — clean baseload stays online longer",
], top=Inches(2.75), left=Inches(0.9), width=Inches(5.7), height=Inches(4.5), size=13)

# Who benefits column
_add_text(s5, Inches(7.0), Inches(2.3), Inches(5.5), Inches(0.4),
          "WHO BENEFITS", size=12, color=RED, bold=True)
_add_bullets(s5, [
    "Plant operators — fewer surprises, better dispatch decisions",
    "Grid operators (ISOs / RTOs) — earlier visibility into capacity risk",
    "Ratepayers — less reliance on costly backup generation during peak demand",
    "The grid as a whole — a more resilient, lower-emission energy system",
], top=Inches(2.75), left=Inches(7.0), width=Inches(5.5), height=Inches(4.5), size=14)

prs.save("slides.pptx")
print("Wrote slides.pptx")
